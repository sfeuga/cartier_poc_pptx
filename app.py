#!/usr/bin/env python
# -*- coding: utf-8 -*-

from flask import Flask, abort, jsonify, request, send_file
from urllib import request, error
from os import environ as env, path
import json

app = Flask(__name__)


STATIC_FOLDER_PATH = "./static"


def get_file_by_name(filename):
    file_path = path.join(STATIC_FOLDER_PATH, filename)

    return file_path


def get_response(url):
    try:
        response = request.urlopen(url)
    except error.URLError as e:
        print('Error:', e.reason)
        return None

    if response and response.getcode() == 200:
        raw_data = response.read()
        data = json.loads(raw_data)
        return data
    else:
        print("Error receiving data", response.getcode())
        return None


@app.route('/api/v1/datas', methods=['GET'])
def get_data():
    mock_api = env.get('mock_api') or 'http://127.0.0.1:8090/api/v1/datas'
    data = get_response(mock_api)

    if data is not None:
        return jsonify(data)
    else:
        abort(42)


@app.route('/api/v1/slides', methods=['GET'])
def generate_slides():
    from pptx import Presentation
    from pptx.util import Cm, Pt
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.dml import MSO_LINE
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # Line Styling
    text_color = RGBColor(127, 127, 127)
    line_color = RGBColor(191, 191, 191)
    dash = MSO_LINE.DASH
    line_width = Pt(0.5)

    prs = Presentation()
    prs.slide_width = Cm(34)
    prs.slide_height = Cm(19)
    blank_slide_layout = prs.slide_layouts[6]

    # Jewel Slide
    slide = prs.slides.add_slide(blank_slide_layout)

    # Title
    y = Cm(0.2)
    x = Cm(2)
    h = Cm(1.5)
    w = Cm(30)
    page_title = slide.shapes.add_textbox(x, y, w, h)
    text_frame = page_title.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    title = paragraph.add_run()
    title.text = 'LOVE'
    title_font = title.font
    title_font.name = 'Brilliant Cut'
    title_font.size = Pt(22)
    title_font.color.rgb = text_color
    sub_title = text_frame.add_paragraph()
    sub_title.text = 'BJ/NJ – 19 MODELS / 20 REFS(1)'
    sub_title.alignment = PP_ALIGN.CENTER
    sub_title.size = Pt(12)
    sub_title.font.bold = True

    # H_LINE_1
    x = Cm(0.7)
    y = Cm(17.7)
    w = Cm(33)
    h = Cm(0)
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LINE_INVERSE, x, y, w, h)
    line = shape.line
    line.dash_style = dash
    line.color.rgb = line_color
    line.width = line_width
    shadow = shape.shadow
    shadow.inherit = False
    # H_LINE_2
    y = Cm(15.5)
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LINE_INVERSE, x, y, w, h)
    line = shape.line
    line.dash_style = dash
    line.color.rgb = line_color
    line.width = line_width
    shadow = shape.shadow
    shadow.inherit = False
    # H_LINE_3
    shapes = slide.shapes
    y = Cm(7)
    shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LINE_INVERSE, x, y, w, h)
    line = shape.line
    line.dash_style = dash
    line.color.rgb = line_color
    line.width = line_width
    shadow = shape.shadow
    shadow.inherit = False
    # H_LINE_4
    shapes = slide.shapes
    y = Cm(4.3)
    shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LINE_INVERSE, x, y, w, h)
    line = shape.line
    line.dash_style = dash
    line.color.rgb = line_color
    line.width = line_width
    shadow = shape.shadow
    shadow.inherit = False
    # V_LINE_1
    shapes = slide.shapes
    x = Cm(8.8)
    y = Cm(2.5)
    w = Cm(0)
    h = Cm(16.4)
    shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LINE_INVERSE, x, y, w, h)
    line = shape.line
    line.dash_style = dash
    line.color.rgb = line_color
    line.width = line_width
    shadow = shape.shadow
    shadow.inherit = False
    # V_LINE_2
    shapes = slide.shapes
    x = Cm(20)
    shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LINE_INVERSE, x, y, w, h)
    line = shape.line
    line.dash_style = dash
    line.color.rgb = line_color
    line.width = line_width
    shadow = shape.shadow
    shadow.inherit = False

    # PRICE_500
    x = Cm(0)
    y = Cm(17.2)
    w = Cm(2.5)
    h = Cm(1)
    text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    price = paragraph.add_run()
    price.text = '€500'
    price_font = price.font
    price_font.name = 'Brilliant Cut Medium'
    price_font.size = Pt(9)
    price_font.color.rgb = text_color
    # PRICE_1000
    y = Cm(14.9)
    text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    price = paragraph.add_run()
    price.text = '€1,000'
    price_font = price.font
    price_font.name = 'Brilliant Cut Medium'
    price_font.size = Pt(9)
    price_font.color.rgb = text_color
    # PRICE_5000
    y = Cm(6.4)
    text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    price = paragraph.add_run()
    price.text = '€5,000'
    price_font = price.font
    price_font.name = 'Brilliant Cut Medium'
    price_font.size = Pt(9)
    price_font.color.rgb = text_color
    # PRICE_10000
    y = Cm(3.7)
    text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    price = paragraph.add_run()
    price.text = '€10,000'
    price_font = price.font
    price_font.name = 'Brilliant Cut Medium'
    price_font.size = Pt(9)
    price_font.color.rgb = text_color

    # W_BANDS_1
    x = Cm(2.57)
    y = Cm(15.8)
    h = Cm(0.41)
    w_band_1_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/W_BAND_1.png', x, y, height=h)
    x = Cm(1.8)
    y = Cm(16.23)
    w = Cm(2.3)
    h = Cm(1.27)
    w_band_1_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = w_band_1_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'B4085200'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'Mini LOVE'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'PG 3g'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '952'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([w_band_1_img, w_band_1_text_box])
    # W_BANDS_2
    x = Cm(2.48)
    y = Cm(11.65)
    h = Cm(0.47)
    w_band_2_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/W_BAND_2.png', x, y, height=h)
    x = Cm(1.7)
    y = Cm(12.14)
    w = Cm(2.48)
    h = Cm(1.28)
    w_band_2_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = w_band_2_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'B4050700'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'Mini LOVE'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'PG 4g 1 DIA 0.02 CT'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '1,804'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([w_band_2_img, w_band_2_text_box])
    # W_BANDS_3
    x = Cm(2.4)
    y = Cm(8.41)
    h = Cm(0.51)
    w_band_3_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/W_BAND_3.png', x, y, height=h)
    x = Cm(1.41)
    y = Cm(8.91)
    w = Cm(2.86)
    h = Cm(1.54)
    w_band_3_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = w_band_3_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'B4050800'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'Mini LOVE'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'PG 4g 8 DIA 0.19 CT'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '3,182'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([w_band_3_img, w_band_3_text_box])
    # W_BANDS_4
    x = Cm(7.18)
    y = Cm(7.22)
    h = Cm(0.45)
    w_band_4_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/W_BAND_4.png', x, y, height=h)
    x = Cm(6.53)
    y = Cm(7.64)
    w = Cm(2.18)
    h = Cm(1.54)
    w_band_4_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = w_band_4_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'B4050600'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'Mini LOVE'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'WG 4g 8 DIA 0.19 CT'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '3,405'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([w_band_4_img, w_band_4_text_box])
    # W_BANDS_5
    x = Cm(2.49)
    y = Cm(5.03)
    h = Cm(0.43)
    w_band_5_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/W_BAND_5.jpg', x, y, height=h)
    x = Cm(1.56)
    y = Cm(5.5)
    w = Cm(2.78)
    h = Cm(1.28)
    w_band_5_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = w_band_5_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'B4085800'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'Mini LOVE'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'PG 3g PAV 0.31 CT'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '5,919'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([w_band_5_img, w_band_5_text_box])

    # RING_1
    x = Cm(12.05)
    y = Cm(13.60)
    h = Cm(0.5)
    ring_1_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/RING_1.png', x, y, height=h)
    x = Cm(11.27)
    y = Cm(14.05)
    w = Cm(2.34)
    ring_1_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = ring_1_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'B4084800'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'Classique'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'PG 5g'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '1,481'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([ring_1_img, ring_1_text_box])
    # RING_2
    x = Cm(9.9)
    y = Cm(12.05)
    h = Cm(0.5)
    ring_2_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/RING_2.png', x, y, height=h)
    x = Cm(9.35)
    y = Cm(12.6)
    w = Cm(1.92)
    ring_2_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = ring_2_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'B4087800'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'PDC'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'PG 9g'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '2,468'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([ring_2_img, ring_2_text_box])
    # RING_3
    x = Cm(13.85)
    y = Cm(11.21)
    h = Cm(1.07)
    ring_3_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/RING_3.png', x, y, height=h)
    x = Cm(14.02)
    y = Cm(12.15)
    w = Cm(1.12)
    ring_3_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = ring_3_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'B4227800'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'LM'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'YG 11g'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '2,898'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([ring_3_img, ring_3_text_box])
    # RING_4
    x = Cm(9.9)
    y = Cm(9.7)
    h = Cm(0.75)
    ring_4_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/RING_4.jpg', x, y, height=h)
    x = Cm(8.95)
    y = Cm(10.54)
    w = Cm(2.63)
    ring_4_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = ring_4_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'B4094300'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'B-LOVE'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'PG WG 6g DIA 0.07 CT'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '3,113'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([ring_4_img, ring_4_text_box])
    # RING_5
    x = Cm(12)
    y = Cm(7.3)
    h = Cm(0.45)
    ring_5_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/RING_5.png', x, y, height=h)
    x = Cm(11.16)
    y = Cm(7.8)
    w = Cm(2.6)
    h = Cm(1.54)
    ring_5_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = ring_5_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'B4218100'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'SM'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'PG 2g PAV 0.19 CT'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '4,011'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([ring_5_img, ring_5_text_box])
    # RING_6
    x = Cm(9.85)
    y = Cm(4.73)
    h = Cm(0.65)
    ring_6_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/RING_6.png', x, y, height=h)
    x = Cm(8.9)
    y = Cm(5.52)
    w = Cm(2.8)
    h = Cm(1.55)
    ring_6_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = ring_6_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'B4094600'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'B-LOVE'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'PG 6g ½ DIA WG PAV 0.19 CT'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '5,750'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([ring_6_img, ring_6_text_box])
    # RING_7
    x = Cm(16.1)
    y = Cm(2.3)
    h = Cm(0.66)
    ring_7_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/RING_7.png', x, y, height=h)
    x = Cm(15.07)
    y = Cm(2.93)
    w = Cm(2.98)
    ring_7_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = ring_7_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'N4210400'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'Classique'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'WG 9g 7 DIA PAV 1.26 CT'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '13,293'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([ring_7_img, ring_7_text_box])

    # NECKPLACE_1
    x = Cm(20.56)
    y = Cm(13.22)
    h = Cm(0.76)
    neckplace_1_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/NECKPLACE_1.png', x, y, height=h)
    x = Cm(20.21)
    y = Cm(14.08)
    w = Cm(2.31)
    neckplace_1_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = neckplace_1_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'B7212300'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'Baby LOVE'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'PG 6g'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '1,931'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([neckplace_1_img, neckplace_1_text_box])
    # NECKPLACE_2
    x = Cm(25.21)
    y = Cm(13.19)
    h = Cm(0.8)
    neckplace_2_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/NECKPLACE_2.jpg', x, y, height=h)
    x = Cm(24.79)
    y = Cm(13.96)
    w = Cm(2.49)
    neckplace_2_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = neckplace_2_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'B7219500'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'Mini Circle'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'YG 6g 2 DIA 0.03 CT'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '1,997'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([neckplace_2_img, neckplace_2_text_box])
    # NECKPLACE_3
    x = Cm(20.78)
    y = Cm(10.39)
    h = Cm(1.21)
    neckplace_3_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/NECKPLACE_3.png', x, y, height=h)
    x = Cm(20.39)
    y = Cm(11.73)
    w = Cm(2.11)
    neckplace_3_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = neckplace_3_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'B7014400'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'Circle of LOVE'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'PG 13g'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '2,708'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([neckplace_3_img, neckplace_3_text_box])
    # NECKPLACE_4
    x = Cm(24.66)
    y = Cm(10.53)
    h = Cm(0.98)
    neckplace_4_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/NECKPLACE_4.png', x, y, height=h)
    x = Cm(24.77)
    y = Cm(11.47)
    w = Cm(1.87)
    h = Cm(1.54)
    neckplace_4_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = neckplace_4_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'B7219700'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'EOL'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'WG PG 8g DIA 0.01 CT'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '2,807'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([neckplace_4_img, neckplace_4_text_box])
    # NECKPLACE_5
    x = Cm(20.57)
    y = Cm(7.86)
    h = Cm(0.9)
    neckplace_5_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/NECKPLACE_5.png', x, y, height=h)
    x = Cm(20.39)
    y = Cm(8.8)
    w = Cm(2.26)
    neckplace_5_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = neckplace_5_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'B7013900'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'Baby LOVE'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'PG 6g PAV 0.22 CT'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '3,755'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([neckplace_5_img, neckplace_5_text_box])
    # NECKPLACE_6
    x = Cm(25.11)
    y = Cm(7.31)
    h = Cm(1.26)
    neckplace_6_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/NECKPLACE_6.png', x, y, height=h)
    x = Cm(24.41)
    y = Cm(8.61)
    w = Cm(2.2)
    neckplace_6_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = neckplace_6_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'B7014600'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'Circle of LOVE'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'WG 14g 3 DIA 0.07 CT'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '3,801'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([neckplace_6_img, neckplace_6_text_box])
    # NECKPLACE_7
    x = Cm(23.5)
    y = Cm(4.36)
    h = Cm(1.11)
    neckplace_7_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/NECKPLACE_7.png', x, y, height=h)
    x = Cm(23.39)
    y = Cm(5.52)
    w = Cm(2.28)
    h = Cm(1.28)
    neckplace_7_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = neckplace_7_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'B7224527'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'Circle of LOVE'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'PG 17g DIA 0.34 CT'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '6,664'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([neckplace_7_img, neckplace_7_text_box])
    # NECKPLACE_8
    x = Cm(22.38)
    y = Cm(2.16)
    h = Cm(0.93)
    neckplace_8_img = slide.shapes.add_picture(STATIC_FOLDER_PATH + '/NECKPLACE_8.png', x, y, height=h)
    x = Cm(22.13)
    y = Cm(3.11)
    w = Cm(1.87)
    neckplace_8_text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = neckplace_8_text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    ref = paragraph.add_run()
    ref.text = 'B7224761-63'
    ref_font = ref.font
    ref_font.name = 'Fancy Cut Medium'
    ref_font.size = Pt(6)
    ref_font.color.rgb = text_color
    mc = text_frame.add_paragraph()
    mc.alignment = PP_ALIGN.CENTER
    mc.text = 'SM'
    mc_font = mc.font
    mc_font.name = 'Fancy Cut Medium'
    mc_font.size = Pt(6)
    mc_font.color.rgb = text_color
    dn = text_frame.add_paragraph()
    dn.alignment = PP_ALIGN.CENTER
    dn.text = 'YG 59g'
    dn_font = dn.font
    dn_font.name = 'Fancy Cut Medium'
    dn_font.size = Pt(6)
    dn_font.color.rgb = RGBColor(248, 120, 214)
    wp = text_frame.add_paragraph()
    wp.alignment = PP_ALIGN.CENTER
    wp.text = '11,837'
    wp_font = wp.font
    wp_font.name = 'Fancy Cut Medium'
    wp_font.size = Pt(6)
    wp_font.color.rgb = text_color
    shapes.add_group_shape([neckplace_8_img, neckplace_8_text_box])

    # FOOTER_1
    x = Cm(3.5)
    y = Cm(17.8)
    w = Cm(3)
    h = Cm(1)
    text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = text_box.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    header = paragraph.add_run()
    header.text = 'WEDDING BANDS'
    header_font = header.font
    header_font.name = 'Brilliant Cut'
    header_font.size = Pt(9)
    header_font.bold = True
    header_font.italic = None
    sub_header = text_frame.add_paragraph()
    sub_header.text = '(5 models / 5 refs)'
    sub_header_font = sub_header.font
    sub_header_font.name = 'Brilliant Cut'
    sub_header_font.size = Pt(9)
    sub_header_font.bold = False
    sub_header_font.italic = None
    # FOOTER_2
    x = Cm(13)
    text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = text_box.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    header = paragraph.add_run()
    header.text = 'RINGS'
    header_font = header.font
    header_font.name = 'Brilliant Cut'
    header_font.size = Pt(9)
    header_font.bold = True
    header_font.italic = None
    sub_header = text_frame.add_paragraph()
    sub_header.text = '(7 models / 7 refs)'
    sub_header_font = sub_header.font
    sub_header_font.name = 'Brilliant Cut'
    sub_header_font.size = Pt(9)
    sub_header_font.bold = False
    sub_header_font.italic = None
    # FOOTER_3
    x = Cm(25.5)
    text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = text_box.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    header = paragraph.add_run()
    header.text = 'NECKLACES'
    header_font = header.font
    header_font.name = 'Brilliant Cut'
    header_font.size = Pt(9)
    header_font.bold = True
    header_font.italic = None
    sub_header = text_frame.add_paragraph()
    sub_header.text = '(7 models / 8 refs)'
    sub_header_font = sub_header.font
    sub_header_font.name = 'Brilliant Cut'
    sub_header_font.size = Pt(9)
    sub_header_font.bold = False
    sub_header_font.italic = None

    # Bags Slide
    slide = prs.slides.add_slide(blank_slide_layout)

    # Title
    y = Cm(0.2)
    x = Cm(2)
    h = Cm(1.5)
    w = Cm(30)
    page_title = slide.shapes.add_textbox(x, y, w, h)
    text_frame = page_title.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    title = paragraph.add_run()
    title.text = 'COLLECTION MODELARIO WOMEN LLG BU23'
    title_font = title.font
    title_font.name = 'Brilliant Cut'
    title_font.size = Pt(22)
    title_font.color.rgb = text_color
    sub_title = text_frame.add_paragraph()
    sub_title.text = '18 MODELS'
    sub_title.alignment = PP_ALIGN.CENTER
    sub_title.size = Pt(12)
    sub_title.font.bold = True

    # H_LINE
    shapes = slide.shapes
    shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LINE_INVERSE, Cm(0.7), Cm(18.15), Cm(33), Cm(0))
    line = shape.line
    line.dash_style = dash
    line.color.rgb = line_color
    line.width = line_width
    shadow = shape.shadow
    shadow.inherit = False
    # V_LINE1
    shapes = slide.shapes
    shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LINE_INVERSE, Cm(7.8), Cm(2.5), Cm(0), Cm(16.4))
    line = shape.line
    line.dash_style = dash
    line.color.rgb = line_color
    line.width = line_width
    shadow = shape.shadow
    shadow.inherit = False
    # V_LINE2
    shapes = slide.shapes
    shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LINE_INVERSE, Cm(14.5), Cm(2.5), Cm(0), Cm(16.4))
    line = shape.line
    line.dash_style = dash
    line.color.rgb = line_color
    line.width = line_width
    shadow = shape.shadow
    shadow.inherit = False
    # V_LINE3
    shapes = slide.shapes
    shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LINE_INVERSE, Cm(20.5), Cm(2.5), Cm(0), Cm(16.4))
    line = shape.line
    line.dash_style = dash
    line.color.rgb = line_color
    line.width = line_width
    shadow = shape.shadow
    shadow.inherit = False
    # V_LINE4
    shapes = slide.shapes
    shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LINE_INVERSE, Cm(27.2), Cm(2.5), Cm(0), Cm(16.4))
    line = shape.line
    line.dash_style = dash
    line.color.rgb = line_color
    line.width = line_width
    shadow = shape.shadow
    shadow.inherit = False

    # IMAGE_BAG_1
    x = Cm(4)
    y = Cm(13.5)
    h = Cm(2.5)
    slide.shapes.add_picture(STATIC_FOLDER_PATH + '/BAG_1.png', x, y, height=h)
    # IMAGE_BAG_2
    x = Cm(8)
    y = Cm(8.5)
    h = Cm(5.5)
    slide.shapes.add_picture(STATIC_FOLDER_PATH + '/BAG_2.png', x, y, height=h)
    # IMAGE_BAG_3
    x = Cm(15)
    y = Cm(5.7)
    h = Cm(7)
    slide.shapes.add_picture(STATIC_FOLDER_PATH + '/BAG_3.png', x, y, height=h)
    # IMAGE_BAG_4
    x = Cm(21)
    y = Cm(3)
    h = Cm(8.5)
    slide.shapes.add_picture(STATIC_FOLDER_PATH + '/BAG_4.png', x, y, height=h)
    # IMAGE_BAG_5
    x = Cm(27.8)
    y = Cm(4.8)
    h = Cm(4.56)
    slide.shapes.add_picture(STATIC_FOLDER_PATH + '/BAG_5.png', x, y, height=h)

    # PRICE_1000
    x = Cm(0)
    y = Cm(17.6)
    h = Cm(1)
    w = Cm(2.5)
    text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    price = paragraph.add_run()
    price.text = '€1,000'
    price_font = price.font
    price_font.name = 'Brilliant Cut Medium'
    price_font.size = Pt(9)
    price_font.color.rgb = text_color
    # PRICE_2500
    y = Cm(9.15)
    text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    price = paragraph.add_run()
    price.text = '€2,500'
    price_font = price.font
    price_font.name = 'Brilliant Cut Medium'
    price_font.size = Pt(9)
    price_font.color.rgb = text_color

    # FOOTER_BAG_1
    x = Cm(2.8)
    y = Cm(18)
    h = Cm(1)
    w = Cm(4)
    text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = text_box.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    header = paragraph.add_run()
    header.text = 'GUIRLANDE DE CARTIER*'
    header_font = header.font
    header_font.name = 'Brilliant Cut'
    header_font.size = Pt(9)
    header_font.bold = True
    header_font.italic = None
    sub_header = text_frame.add_paragraph()
    sub_header.text = '1 MODEL'
    sub_header_font = sub_header.font
    sub_header_font.name = 'Brilliant Cut'
    sub_header_font.size = Pt(9)
    sub_header_font.bold = False
    sub_header_font.italic = None
    # FOOTER_BAG_2
    x = Cm(10)
    w = Cm(2)
    text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = text_box.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    header = paragraph.add_run()
    header.text = 'MUST*'
    header_font = header.font
    header_font.name = 'Brilliant Cut'
    header_font.size = Pt(9)
    header_font.bold = True
    header_font.italic = None
    sub_header = text_frame.add_paragraph()
    sub_header.text = '5 MODELS'
    sub_header_font = sub_header.font
    sub_header_font.name = 'Brilliant Cut'
    sub_header_font.size = Pt(9)
    sub_header_font.bold = False
    sub_header_font.italic = None
    # FOOTER_BAG_3
    x = Cm(16)
    w = Cm(3.5)
    text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = text_box.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    header = paragraph.add_run()
    header.text = 'DOUBLE C DE CARTIER*'
    header_font = header.font
    header_font.name = 'Brilliant Cut'
    header_font.size = Pt(9)
    header_font.bold = True
    header_font.italic = None
    sub_header = text_frame.add_paragraph()
    sub_header.text = '4 MODELS'
    sub_header_font = sub_header.font
    sub_header_font.name = 'Brilliant Cut'
    sub_header_font.size = Pt(9)
    sub_header_font.bold = False
    sub_header_font.italic = None
    # FOOTER_BAG_4
    x = Cm(23)
    w = Cm(2.5)
    text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = text_box.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    header = paragraph.add_run()
    header.text = 'PANTHERE C*'
    header_font = header.font
    header_font.name = 'Brilliant Cut'
    header_font.size = Pt(9)
    header_font.bold = True
    header_font.italic = None
    sub_header = text_frame.add_paragraph()
    sub_header.text = '5 MODELS'
    sub_header_font = sub_header.font
    sub_header_font.name = 'Brilliant Cut'
    sub_header_font.size = Pt(9)
    sub_header_font.bold = False
    sub_header_font.italic = None
    # FOOTER_BAG_5
    x = Cm(29)
    w = Cm(3.3)
    text_box = slide.shapes.add_textbox(x, y, w, h)
    text_frame = text_box.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    header = paragraph.add_run()
    header.text = 'PANTHERE HEAD*'
    header_font = header.font
    header_font.name = 'Brilliant Cut'
    header_font.size = Pt(9)
    header_font.bold = True
    header_font.italic = None
    sub_header = text_frame.add_paragraph()
    sub_header.text = '3 MODELS'
    sub_header_font = sub_header.font
    sub_header_font.name = 'Brilliant Cut'
    sub_header_font.size = Pt(9)
    sub_header_font.bold = False
    sub_header_font.italic = None

    prs.save(STATIC_FOLDER_PATH + '/slides.pptx')

    return send_file(STATIC_FOLDER_PATH + '/slides.pptx',
                     mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                     as_attachment=True,
                     download_name='slides.pptx',
                     attachment_filename='slides.pptx')


if __name__ == '__main__':
    app.run()
